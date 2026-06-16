# -*- coding: utf-8 -*-
"""
学生校内双证管理系统 - 路演 PPT 生成脚本（管理端 + 学生端双平台）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
PRIMARY = RGBColor(0x16, 0x77, 0xFF)      # 主题蓝
DARK = RGBColor(0x00, 0x1C, 0x40)         # 深色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF5, 0xFF)     # 浅蓝背景
ACCENT_GREEN = RGBColor(0x52, 0xC4, 0x1A)
ACCENT_ORANGE = RGBColor(0xFA, 0x8C, 0x16)
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4F)
ACCENT_PURPLE = RGBColor(0x72, 0x2E, 0xD1)
ACCENT_PINK = RGBColor(0xEB, 0x2F, 0x96)
GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DARK_TEXT = RGBColor(0x26, 0x26, 0x26)
CARD_BG = RGBColor(0xE6, 0xF7, 0xFF)
CARD_BG_GREEN = RGBColor(0xF0, 0xFF, 0xF0)
STUDENT_BG = RGBColor(0xF5, 0xF0, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)


# ============================================================
# Helper functions
# ============================================================
def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    return shape

def add_rounded_rect(slide, left, top, width, height, color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT,
                   bold_first=False, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(font_size * (line_spacing - 1) * 0.7)
    return txBox

def add_section_header(slide, number, title, subtitle=""):
    """左侧深色数字 + 标题区域"""
    add_rect(slide, Inches(0), Inches(0), Inches(4.5), Inches(7.5), color=PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(1.5),
                 f"0{number}", font_size=72, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(3.2), Inches(1),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(4.5), Inches(3.2), Inches(0.6),
                     subtitle, font_size=14, color=RGBColor(0xBB, 0xCC, 0xFF))


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_rect(slide, Inches(1.5), Inches(2.0), Inches(2.5), Inches(0.015), color=PRIMARY)
add_rect(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.008), color=RGBColor(0x33, 0x55, 0x88))

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
             '学生校内双证管理系统', font_size=46, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.05), Inches(10), Inches(0.5),
             'Dual Certificate Management System for Students on Campus', font_size=18, color=RGBColor(0x99, 0xBB, 0xDD))
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.5),
             '管理端 + 学生端   |   人社/专业/校内引进三类证书   |   报名→考试→AI审核→归档  全流程管理',
             font_size=13, color=GRAY)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
             '技术栈：React 18 + TypeScript + Ant Design 5   |   Node.js + Express + SQLite   |   JWT 双角色认证   |   AI 智能审核引擎',
             font_size=12, color=GRAY)
add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.4),
             '管理端 :5173 · 学生端 :5174 · 公开门户 /public · 后端 API :3001',
             font_size=11, color=RGBColor(0x66, 0x88, 0xAA))


# ============================================================
# Slide 2: 项目背景与痛点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 1, '项目背景', 'Project Background')

# 痛点卡片
shape = add_rect(slide, Inches(5.2), Inches(0.8), Inches(7.3), Inches(2.8), color=LIGHT_BG)
add_text_box(slide, Inches(5.5), Inches(1.0), Inches(6.8), Inches(0.4),
             '🎯 核心痛点', font_size=18, color=PRIMARY, bold=True)
add_multi_text(slide, Inches(5.5), Inches(1.5), Inches(6.8), Inches(2.0), [
    '• 学校内部多种证书管理混乱 — 人社、专业、校内引进三种类型各自为政',
    '• 报名→考试→取证全流程依赖纸质表格和人工传递，效率低下',
    '• 成绩审核缺乏统一标准，人工审核主观性强、工作量巨大',
    '• 证书数据分散在多个系统，无法统一查询统计和班级归档',
    '• 学生缺乏自助查询入口，培训信息和辅导材料缺少统一公示平台',
    '• 管理员和学生之间信息不对称，学生无法实时了解审核进度',
], font_size=12, line_spacing=1.8)

# 目标卡片
shape = add_rect(slide, Inches(5.2), Inches(4.0), Inches(7.3), Inches(3.0), color=CARD_BG)
add_text_box(slide, Inches(5.5), Inches(4.2), Inches(6.8), Inches(0.4),
             '✅ 系统目标', font_size=18, color=ACCENT_GREEN, bold=True)
add_multi_text(slide, Inches(5.5), Inches(4.7), Inches(6.8), Inches(2.2), [
    '• 建立统一的校内双证管理平台，覆盖三种证书类型的全生命周期管理',
    '• 管理端 + 学生端双平台架构，管理员全局管控，学生自助操作',
    '• 报名→审核→考试→AI审核→归档的全流程数字化闭环',
    '• 通过 AI 辅助审核提升审核效率与标准一致性',
    '• 按班级自动归档，一键生成 Excel 报表',
    '• 公开门户实现证书信息/培训/材料的统一公示',
], font_size=12, line_spacing=1.8)


# ============================================================
# Slide 3: 系统架构（双平台）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 2, '系统架构', 'System Architecture')

# 双平台展示
# 管理端
add_rounded_rect(slide, Inches(5.3), Inches(0.5), Inches(3.5), Inches(2.8), color=LIGHT_BG)
add_text_box(slide, Inches(5.5), Inches(0.6), Inches(3.0), Inches(0.4),
             '🔧 管理端 (Port 5173)', font_size=15, color=PRIMARY, bold=True)
add_text_box(slide, Inches(5.5), Inches(0.95), Inches(3.0), Inches(0.3),
             'admin / admin123', font_size=10, color=GRAY)
add_multi_text(slide, Inches(5.5), Inches(1.3), Inches(3.0), Inches(1.8), [
    '📊 仪表盘 — 多维度数据统计',
    '👨‍🎓 学生管理 — CRUD + 批量导入',
    '📜 证书管理 — 三类证书 + 报名规则',
    '📝 报名管理 — 报名审核（批准/拒绝）',
    '📄 考试管理 — AI审核 + 人工复核',
    '🏆 证书记录 — 获取情况管理',
    '📦 归档管理 — 按班归档 + Excel导出',
    '📖 培训 & 辅导材料管理',
], font_size=9, line_spacing=1.5)

# 学生端
add_rounded_rect(slide, Inches(9.2), Inches(0.5), Inches(3.5), Inches(2.8), color=STUDENT_BG)
add_text_box(slide, Inches(9.4), Inches(0.6), Inches(3.0), Inches(0.4),
             '🎓 学生端 (Port 5174)', font_size=15, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(9.4), Inches(0.95), Inches(3.0), Inches(0.3),
             '学号 / 密码（学号后6位）', font_size=10, color=GRAY)
add_multi_text(slide, Inches(9.4), Inches(1.3), Inches(3.0), Inches(1.8), [
    '📊 我的仪表盘 — 个人统计',
    '🎯 证书报名 — 浏览 + 一键报名',
    '📋 我的报名 — 审核状态查看',
    '📝 我的考试 — 考试成绩提交',
    '🏆 我的证书 — 已获证书记录',
    '📖 学习资源 — 培训 + 辅导材料',
    '👤 个人信息 — 查看 & 修改密码',
], font_size=9, line_spacing=1.5)

add_text_box(slide, Inches(8.8), Inches(1.5), Inches(0.4), Inches(0.4),
             '⚡', font_size=24, alignment=PP_ALIGN.CENTER)

# 连接箭头
add_text_box(slide, Inches(6.8), Inches(3.3), Inches(0.5), Inches(0.4),
             '▼', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.5), Inches(3.3), Inches(0.5), Inches(0.4),
             '▼', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# 后端层
add_rounded_rect(slide, Inches(5.3), Inches(3.7), Inches(7.4), Inches(1.5), color=RGBColor(0x08, 0x99, 0xFF))
add_text_box(slide, Inches(5.5), Inches(3.8), Inches(7), Inches(0.35),
             '⚙️ Express Server (Port 3001) — Node.js + TypeScript', font_size=14, color=WHITE, bold=True)
add_multi_text(slide, Inches(5.5), Inches(4.2), Inches(7), Inches(0.8), [
    'JWT 双角色认证 (admin + student)  |  10 个 Controller 模块  |  AI 审核引擎  |  Multer 文件上传  |  XLSX 解析',
    '50+ RESTful API 端点 — 管理端 30+  /  学生端 10+  /  公开接口 2  /  教务对接 2',
], font_size=10, color=RGBColor(0xDD, 0xEE, 0xFF), line_spacing=1.4)

add_text_box(slide, Inches(8.5), Inches(5.2), Inches(0.5), Inches(0.4),
             '▼', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# 数据层
add_rounded_rect(slide, Inches(5.3), Inches(5.6), Inches(7.4), Inches(1.0), color=DARK)
add_text_box(slide, Inches(5.5), Inches(5.65), Inches(7), Inches(0.3),
             '🗄️ SQLite (sql.js — 纯 JS 实现，零安装依赖)', font_size=14, color=WHITE, bold=True)
add_text_box(slide, Inches(5.5), Inches(6.05), Inches(7), Inches(0.4),
             'users | students(含密码) | certificate_types | certificates | registration_rules | training_materials | student_registrations | exam_submissions | certificate_records | archives | archive_details',
             font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC))

# 公开门户
add_rounded_rect(slide, Inches(5.3), Inches(6.9), Inches(7.4), Inches(0.4), color=ACCENT_GREEN)
add_text_box(slide, Inches(5.5), Inches(6.93), Inches(7), Inches(0.3),
             '🌐 公开门户 /public — 证书信息展示 + 培训公示 + 辅导材料查阅（无需登录）',
             font_size=10, color=WHITE, bold=True)


# ============================================================
# Slide 4: 证书类型体系
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 3, '证书类型体系', 'Certificate Categories')

cert_types = [
    ('1', '人社证书', 'RENSHE', '人社部门颁发的\n职业资格证书\n（必须获取）', PRIMARY, [
        '计算机等级考试二级', '英语四级(CET-4)', '电工证', '焊工证'
    ]),
    ('2', '专业证书', 'ZHUANYE', '专业相关的\n行业认证证书\n（必须获取）', ACCENT_GREEN, [
        '会计从业资格证', '软件设计师', '建造师资格证', '网络工程师'
    ]),
    ('3', '校内引进证书', 'XIAONEI', '学校引进的\n技能认证证书\n（可选获取）', ACCENT_ORANGE, [
        '普通话等级证书', '电子商务师', '市场营销师', '数据分析师'
    ]),
]

for i, (num, name, code, desc, color, examples) in enumerate(cert_types):
    left = Inches(5.3) + Inches(2.7) * i
    shape = add_rect(slide, left, Inches(1.0), Inches(2.4), Inches(5.8), color=LIGHT_BG)
    add_rect(slide, left, Inches(1.0), Inches(2.4), Inches(0.09), color=color)
    add_text_box(slide, left + Inches(0.2), Inches(1.2), Inches(0.6), Inches(0.6),
                 num, font_size=36, color=color, bold=True)
    add_text_box(slide, left + Inches(0.8), Inches(1.35), Inches(1.4), Inches(0.4),
                 name, font_size=20, color=color, bold=True)
    add_text_box(slide, left + Inches(0.8), Inches(1.8), Inches(1.4), Inches(0.3),
                 code, font_size=10, color=GRAY)
    add_text_box(slide, left + Inches(0.3), Inches(2.4), Inches(1.8), Inches(1.1),
                 desc, font_size=11, color=DARK_TEXT)
    add_text_box(slide, left + Inches(0.3), Inches(3.8), Inches(1.8), Inches(0.3),
                 '示例证书：', font_size=10, color=GRAY, bold=True)
    for j, ex in enumerate(examples):
        add_text_box(slide, left + Inches(0.3), Inches(4.2) + Inches(0.35) * j,
                     Inches(1.8), Inches(0.3), f'▸ {ex}', font_size=10, color=DARK_TEXT)

shape = add_rect(slide, Inches(5.3), Inches(7.0), Inches(7.5), Inches(0.4), color=DARK)
add_text_box(slide, Inches(5.5), Inches(7.03), Inches(7), Inches(0.3),
             '系统支持自定义证书类型和证书信息，管理员可灵活配置报名规则（时间范围、容量限制、报名条件 JSON）',
             font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC))


# ============================================================
# Slide 5: 管理端核心功能
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color=PRIMARY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(5), Inches(0.5),
             '管理端 — 核心功能模块', font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(0.7), Inches(5), Inches(0.4),
             'Admin Portal — 全局数据管理', font_size=14, color=RGBColor(0xBB, 0xDD, 0xFF))

modules = [
    ('📊', '仪表盘', '学生总数/获证总数\n报名总数/待审核数\n类型统计/班级统计\n月度趋势图'),
    ('👨‍🎓', '学生管理', 'CRUD管理\nExcel批量导入\n按班级/学号搜索\n新增自动设密码'),
    ('📜', '证书管理', '3种证书类型维护\n证书信息+颁发机构\n报名规则配置\n时间/容量/条件'),
    ('📝', '报名管理', '学生报名创建\n审核批准/拒绝\n审核意见填写\n批量导入报名'),
    ('📄', '考试管理', '考试成绩录入\nAI智能审核\n人工复核审批\n批量导入考试'),
    ('🏆', '证书记录', '获取情况管理\n证书编号/成绩\n批量导入记录\n自动统计报表'),
    ('📦', '归档管理', '按班级归档\n学生通过率统计\nExcel一键导出\n归档详情查看'),
    ('📖', '培训&材料', '培训信息/辅导材料\n视频/文档/链接\n公开/内部可见\n文件上传'),
]

for i, (icon, title, desc) in enumerate(modules):
    row = i // 4
    col = i % 4
    left = Inches(0.4) + Inches(3.2) * col
    top = Inches(1.5) + Inches(2.85) * row
    shape = add_rect(slide, left, top, Inches(2.95), Inches(2.6), color=CARD_BG)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=30)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.15), Inches(2.0), Inches(0.45),
                 title, font_size=18, color=PRIMARY, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), Inches(2.5), Inches(1.6),
                 desc, font_size=10, color=DARK_TEXT)


# ============================================================
# Slide 6: 学生端核心功能
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color=ACCENT_PURPLE)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(5), Inches(0.5),
             '学生端 — 核心功能模块', font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(0.7), Inches(5), Inches(0.4),
             'Student Portal — 自助报名 & 成绩提交 & 证书查看', font_size=14, color=RGBColor(0xDD, 0xCC, 0xFF))

modules = [
    ('🔐', '学号登录', '学号+密码认证\n默认密码=学号后6位\n首次登录提醒改密\nJWT Token 24h有效'),
    ('📊', '我的仪表盘', '报名数/考试数\n获证数/待审核数\n各类型获取统计\n最近报名记录'),
    ('🎯', '证书报名', '浏览可报名证书\n一键报名按钮\n自动校验名额\n实时显示报名状态'),
    ('📋', '我的报名', '个人全部报名记录\n审核状态查看\n通过/拒绝/审核中\n审核意见显示'),
    ('📝', '我的考试', '考试记录查看\n提交考试成绩\n上传成绩单附件\nAI+人工审核进度'),
    ('🏆', '我的证书', '已获证书记录\n证书编号/成绩\n获取日期展示\n按类型分类查看'),
    ('📖', '学习资源', '培训信息查看\n辅导材料查阅\n视频/文档/链接\nTab分类切换'),
    ('👤', '个人中心', '个人信息查看\n学号/班级/专业\n修改密码功能\n安全退出登录'),
]

for i, (icon, title, desc) in enumerate(modules):
    row = i // 4
    col = i % 4
    left = Inches(0.4) + Inches(3.2) * col
    top = Inches(1.5) + Inches(2.85) * row
    shape = add_rect(slide, left, top, Inches(2.95), Inches(2.6), color=STUDENT_BG)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=30)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.15), Inches(2.0), Inches(0.45),
                 title, font_size=18, color=ACCENT_PURPLE, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), Inches(2.5), Inches(1.6),
                 desc, font_size=10, color=DARK_TEXT)


# ============================================================
# Slide 7: AI 审核引擎
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             '🤖 AI 智能审核引擎', font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
             '三维度智能评估 — 材料完整性 × 成绩合格性 × 前置条件审核', font_size=14, color=RGBColor(0x88, 0xAA, 0xCC))

dimensions = [
    ('📋', '材料完整性', '检查是否上传了\n考试结果附件材料\n（成绩单/证书扫描件）', PRIMARY),
    ('📈', '成绩合格性', '验证考试成绩是否\n达到及格线标准\n（默认 60 分）', ACCENT_GREEN),
    ('✅', '前置条件审核', '确认报名是否已通过\n管理员审核\n确保流程合规性', ACCENT_ORANGE),
]

for i, (icon, title, desc, color) in enumerate(dimensions):
    left = Inches(0.8) + Inches(4.2) * i
    top = Inches(1.8)
    shape = add_rect(slide, left, top, Inches(3.8), Inches(2.5), color=RGBColor(0x0A, 0x2A, 0x50))
    add_rect(slide, left, top, Inches(3.8), Inches(0.06), color=color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=36)
    add_text_box(slide, left + Inches(1.0), top + Inches(0.35), Inches(2.5), Inches(0.45),
                 title, font_size=22, color=color, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.1), Inches(3.2), Inches(1.2),
                 desc, font_size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

# 审核结果
results = [
    ('✅ 通过', '材料完整\n成绩合格\n前置满足', ACCENT_GREEN),
    ('⚠️ 需补充材料', '材料不完整\n需补充附件后重审', ACCENT_ORANGE),
    ('❌ 不通过', '成绩不达标\n或前置条件未满足', ACCENT_RED),
]

add_text_box(slide, Inches(0.8), Inches(4.6), Inches(8), Inches(0.5),
             '审核结果输出', font_size=18, color=WHITE, bold=True)

for i, (title, desc, color) in enumerate(results):
    left = Inches(0.8) + Inches(4.2) * i
    top = Inches(5.2)
    shape = add_rect(slide, left, top, Inches(3.8), Inches(1.5), color=RGBColor(0x0A, 0x2A, 0x50))
    add_rect(slide, left, top, Inches(3.8), Inches(0.06), color=color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.2), Inches(0.45),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.2), Inches(0.7),
                 desc, font_size=11, color=GRAY)

add_text_box(slide, Inches(0.8), Inches(7.0), Inches(12), Inches(0.3),
             '💡 内置模拟AI引擎 (server/src/services/aiReviewService.ts)，审核附置信度评分，可替换为真实 AI API（GPT/Claude）',
             font_size=10, color=RGBColor(0x66, 0x88, 0xAA))


# ============================================================
# Slide 8: 双端业务流程
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 4, '业务流程', 'Business Workflow')

# 管理端流程
add_rounded_rect(slide, Inches(5.2), Inches(0.6), Inches(7.5), Inches(3.0), color=LIGHT_BG)
add_text_box(slide, Inches(5.5), Inches(0.7), Inches(7), Inches(0.35),
             '🔧 管理端 — 全流程管理', font_size=16, color=PRIMARY, bold=True)

admin_steps = [
    ('①', '创建学生', '录入信息/批量导入\n自动设密码(学号后6位)'),
    ('②', '配置证书', '定义证书类型\n设置报名规则'),
    ('③', '审核报名', '批准/拒绝\n填写审核意见'),
    ('④', '管理考试', 'AI+人工审核\n通过后自动生成证书记录'),
    ('⑤', '按班归档', '自动汇总\n一键导出Excel'),
]

for i, (num, title, desc) in enumerate(admin_steps):
    left = Inches(5.4) + Inches(1.5) * i
    top = Inches(1.15)
    add_rounded_rect(slide, left, top, Inches(1.3), Inches(2.0), color=WHITE)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.1), Inches(1.2), Inches(0.3),
                 num, font_size=14, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.4), Inches(1.2), Inches(0.3),
                 title, font_size=11, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.8), Inches(1.2), Inches(1.0),
                 desc, font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, left + Inches(1.3), top + Inches(0.8), Inches(0.2), Inches(0.3),
                     '→', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# 学生端流程
add_rounded_rect(slide, Inches(5.2), Inches(3.9), Inches(7.5), Inches(3.0), color=STUDENT_BG)
add_text_box(slide, Inches(5.5), Inches(4.0), Inches(7), Inches(0.35),
             '🎓 学生端 — 自助操作', font_size=16, color=ACCENT_PURPLE, bold=True)

student_steps = [
    ('①', '学号登录', '学号+密码\n(默认学号后6位)'),
    ('②', '浏览证书', '查看可报名证书\n一键报名'),
    ('③', '查看审核', '报名状态跟踪\n查看审核意见'),
    ('④', '提交考试', '上传成绩单\n填写分数'),
    ('⑤', '查看证书', '已获证书记录\n证书编号+成绩'),
]

for i, (num, title, desc) in enumerate(student_steps):
    left = Inches(5.4) + Inches(1.5) * i
    top = Inches(4.45)
    add_rounded_rect(slide, left, top, Inches(1.3), Inches(2.0), color=WHITE)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.1), Inches(1.2), Inches(0.3),
                 num, font_size=14, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.4), Inches(1.2), Inches(0.3),
                 title, font_size=11, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.8), Inches(1.2), Inches(1.0),
                 desc, font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, left + Inches(1.3), top + Inches(1.1), Inches(0.2), Inches(0.3),
                     '→', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide 9: 认证与权限体系
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 5, '认证与权限', 'Authentication & Authorization')

# JWT 说明
add_rounded_rect(slide, Inches(5.3), Inches(0.6), Inches(7.5), Inches(2.0), color=LIGHT_BG)
add_text_box(slide, Inches(5.5), Inches(0.7), Inches(7), Inches(0.35),
             '🔐 JWT 双角色认证体系', font_size=18, color=PRIMARY, bold=True)
add_multi_text(slide, Inches(5.5), Inches(1.2), Inches(7), Inches(1.2), [
    '• Token 过期时间: 24小时，过期自动跳转登录页',
    '• bcrypt 加密存储密码，默认密码为学号后6位',
    '• 管理端中间件 requireAuth  + 学生端中间件 requireStudentAuth',
    '• 请求拦截器自动注入 Authorization: Bearer <token>',
    '• 401 自动清除 token 并重定向至登录页',
], font_size=11, line_spacing=1.8)

# 双账号体系对比
add_text_box(slide, Inches(5.5), Inches(2.9), Inches(7), Inches(0.4),
             '双端账号体系对比', font_size=16, color=DARK_TEXT, bold=True)

account_data = [
    ('维度', '管理端', '学生端'),
    ('登录方式', '用户名 + 密码', '学号 + 密码'),
    ('登录接口', 'POST /api/auth/login', 'POST /api/student-auth/login'),
    ('Token Role', 'admin / super_admin', 'student'),
    ('密码修改', '--', 'POST /api/student-auth/change-password'),
    ('默认密码', 'admin123', '学号后6位 (bcrypt哈希)'),
    ('权限范围', '全部数据的 CRUD + 审核', '仅自己的报名/考试/证书'),
    ('端口', 'http://localhost:5173', 'http://localhost:5174'),
]

table_top = Inches(3.4)
row_height = Inches(0.4)
col_widths = [Inches(1.8), Inches(2.7), Inches(2.7)]
headers = account_data[0]
# Header row
add_rounded_rect(slide, Inches(5.5), table_top, col_widths[0], row_height, color=PRIMARY)
add_text_box(slide, Inches(5.6), table_top + Inches(0.05), col_widths[0] - Inches(0.2), Inches(0.3),
             headers[0], font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, Inches(5.5) + col_widths[0], table_top, col_widths[1], row_height, color=PRIMARY)
add_text_box(slide, Inches(5.6) + col_widths[0], table_top + Inches(0.05), col_widths[1] - Inches(0.2), Inches(0.3),
             headers[1], font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, Inches(5.5) + col_widths[0] + col_widths[1], table_top, col_widths[2], row_height, color=ACCENT_PURPLE)
add_text_box(slide, Inches(5.6) + col_widths[0] + col_widths[1], table_top + Inches(0.05), col_widths[2] - Inches(0.2), Inches(0.3),
             headers[2], font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for j, row in enumerate(account_data[1:]):
    y = table_top + row_height + row_height * j
    bg_color = CARD_BG if j % 2 == 0 else WHITE
    add_rounded_rect(slide, Inches(5.5), y, col_widths[0], row_height, color=bg_color)
    add_text_box(slide, Inches(5.6), y + Inches(0.05), col_widths[0] - Inches(0.2), Inches(0.3),
                 row[0], font_size=10, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, Inches(5.5) + col_widths[0], y, col_widths[1], row_height, color=bg_color)
    add_text_box(slide, Inches(5.6) + col_widths[0], y + Inches(0.05), col_widths[1] - Inches(0.2), Inches(0.3),
                 row[1], font_size=10, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, Inches(5.5) + col_widths[0] + col_widths[1], y, col_widths[2], row_height, color=bg_color)
    add_text_box(slide, Inches(5.6) + col_widths[0] + col_widths[1], y + Inches(0.05), col_widths[2] - Inches(0.2), Inches(0.3),
                 row[2], font_size=10, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# 默认账号速查
add_text_box(slide, Inches(5.5), Inches(6.3), Inches(7), Inches(0.35),
             '🔑 默认账号速查', font_size=16, color=DARK_TEXT, bold=True)
accounts = [
    '管理员: admin / admin123      (管理端 http://localhost:5173)',
    '学生张三: 2024001 / 024001     (学生端 http://localhost:5174)',
    '学生李四: 2024002 / 24002      (学生端 http://localhost:5174)',
    '学生王五: 2024003 / 24003      (学生端 http://localhost:5174)',
]
for i, acc in enumerate(accounts):
    add_text_box(slide, Inches(5.7), Inches(6.65) + Inches(0.22) * i, Inches(6.5), Inches(0.2),
                 acc, font_size=10, color=DARK_TEXT)


# ============================================================
# Slide 10: API 接口总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 6, 'API 接口', 'RESTful API (50+ Endpoints)')

# 管理端 API
add_text_box(slide, Inches(5.3), Inches(0.5), Inches(7), Inches(0.4),
             '🔧 管理端 API (30+)', font_size=15, color=PRIMARY, bold=True)
api_groups = [
    ('🔐 认证', [('POST /api/auth/login', '管理员登录'), ('GET /api/auth/me', '当前用户')]),
    ('👨‍🎓 学生', [('GET/POST /api/students', 'CRUD + 分页搜索'), ('PUT/DELETE /:id', '编辑删除'), ('POST batch-import', 'Excel批量导入')]),
    ('📜 证书', [('GET/POST /api/certificates', 'CRUD'), ('GET/POST types', '证书类型')]),
    ('📋 规则', [('GET/POST /api/registration-rules', '报名规则')]),
    ('📝 报名', [('GET/POST /api/registrations', '报名CRUD'), ('PUT /:id/review', '审核')]),
    ('📄 考试', [('GET/POST /api/exam-submissions', '考试CRUD'), ('POST /:id/ai-review', 'AI审核'), ('PUT /:id/review', '人工审核')]),
    ('🏆 记录', [('GET/POST /api/certificate-records', 'CRUD'), ('GET /api/stats', '仪表盘统计')]),
    ('📦 归档', [('GET/POST /api/archives', '归档管理'), ('GET /:id/export', '导出Excel')]),
    ('📖 培训', [('GET/POST /api/training-materials', '培训+材料CRUD')]),
    ('🔗 教务', [('GET /api/edusys/certificates/:no', '按学号查询'), ('POST /api/edusys/sync', '批量同步')]),
]
for i, (group, endpoints) in enumerate(api_groups):
    col = i % 5
    row = i // 5
    left = Inches(5.2) + Inches(1.65) * col
    top = Inches(1.0) + Inches(2.5) * row
    shape = add_rect(slide, left, top, Inches(1.5), Inches(2.2), color=CARD_BG)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.05), Inches(1.3), Inches(0.3),
                 group, font_size=11, color=PRIMARY, bold=True)
    for j, (ep, desc) in enumerate(endpoints):
        y = top + Inches(0.35) + Inches(0.55) * j
        add_text_box(slide, left + Inches(0.1), y, Inches(1.3), Inches(0.2),
                     ep, font_size=7, color=PRIMARY, bold=True)
        add_text_box(slide, left + Inches(0.1), y + Inches(0.18), Inches(1.3), Inches(0.2),
                     desc, font_size=7, color=GRAY)

# 学生端 API (右侧下半部)
add_text_box(slide, Inches(5.3), Inches(5.2), Inches(7), Inches(0.35),
             '🎓 学生端 API (12)            🌐 公开接口 (2)            🔗 教务对接 (2)',
             font_size=12, color=ACCENT_PURPLE, bold=True)

shape = add_rect(slide, Inches(5.3), Inches(5.6), Inches(7.5), Inches(1.6), color=STUDENT_BG)
student_apis = [
    ('POST', '/api/student-auth/login', '学号密码登录'),
    ('POST', '/api/student-auth/change-password', '修改密码'),
    ('GET', '/api/student/me', '个人信息'),
    ('GET', '/api/student/dashboard', '个人仪表盘'),
    ('GET', '/api/student/certificates', '可报名证书列表'),
    ('GET', '/api/student/registrations', '我的报名'),
    ('POST', '/api/student/registrations', '一键报名'),
    ('GET', '/api/student/exams', '我的考试'),
    ('POST', '/api/student/exams', '提交考试成绩 (含文件上传)'),
    ('GET', '/api/student/records', '我的证书记录'),
    ('GET', '/api/student/materials', '学习资源'),
    ('GET', '/api/student/approved-certificates', '已通过报名的证书 (用于考试提交下拉)'),
]
for i, (method, path, desc) in enumerate(student_apis):
    col = i // 6
    row = i % 6
    left = Inches(5.4) + Inches(3.7) * col
    top = Inches(5.9) + Inches(0.25) * row
    method_color = ACCENT_GREEN if method == 'GET' else (ACCENT_ORANGE if method == 'POST' else PRIMARY)
    add_text_box(slide, left, top, Inches(0.35), Inches(0.2),
                 method, font_size=7, color=method_color, bold=True)
    add_text_box(slide, left + Inches(0.38), top, Inches(2.5), Inches(0.2),
                 f'{path} — {desc}', font_size=7, color=DARK_TEXT)


# ============================================================
# Slide 11: 技术亮点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 7, '技术亮点', 'Technical Highlights')

highlights = [
    ('🤖', 'AI 智能审核', '内置三维度 AI 审核引擎，自动评估材料完整性、成绩合格性和前置条件。结果附带置信度和详细评语。预留真实 AI API 接口，可对接 GPT/Claude。'),
    ('🔐', 'JWT 双角色认证', '管理端 admin/super_admin + 学生端 student 双角色体系。bcrypt 密码哈希，Token 24h 有效，中间件拦截未授权请求，401 自动跳转登录。'),
    ('⚡', '零配置数据库', '使用 sql.js（纯 JS SQLite），npm install 即用，无需安装任何数据库软件。数据自动持久化到 data.db 文件。'),
    ('📦', '全栈 TypeScript', '前后端统一 TypeScript，编译时类型检查减少运行时错误，IDE 智能提示提升开发效率。管理端和学生端共享技术栈。'),
    ('🎨', '双端独立架构', '管理端 (5173) + 学生端 (5174) 两个独立 Vite 项目，共享同一后端 API。各自有独立的登录页、布局、路由和权限控制。'),
    ('📊', '批量数据处理', '全流程支持 Excel 批量导入（学生/报名/考试/证书记录）。按班级自动归档，一键导出 Excel 报表。xlsx 库多格式兼容。'),
    ('🌐', '信息公示门户', '/public 公开页面，无需登录即可查阅证书信息、培训和辅导材料。可嵌入学校官网或公众号作为信息公示页。'),
    ('🔗', '教务系统对接', '标准 RESTful API 对接教务系统，支持按学号查询和批量同步。模拟数据演示，可快速替换为真实接口。'),
    ('🔄', '报名→归档闭环', '学生报名 → 管理员审核 → 学生提交考试 → AI审核 → 人工复核 → 自动生成证书记录 → 按班归档，完整数字化闭环。'),
]

for i, (icon, title, desc) in enumerate(highlights):
    row = i // 3
    col = i % 3
    left = Inches(5.2) + Inches(2.7) * col
    top = Inches(0.6) + Inches(2.25) * row
    shape = add_rect(slide, left, top, Inches(2.5), Inches(2.0), color=CARD_BG)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.1), Inches(0.4), Inches(0.4),
                 icon, font_size=20)
    add_text_box(slide, left + Inches(0.55), top + Inches(0.1), Inches(1.8), Inches(0.3),
                 title, font_size=14, color=PRIMARY, bold=True)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.5), Inches(2.2), Inches(1.4),
                 desc, font_size=8, color=DARK_TEXT)


# ============================================================
# Slide 12: 项目总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, 8, '项目总结', 'Project Summary')

add_text_box(slide, Inches(5.3), Inches(0.6), Inches(7), Inches(0.45),
             '✅ 系统实现的核心价值', font_size=20, color=PRIMARY, bold=True)

achievements = [
    '覆盖人社/专业/校内引进三类证书的全生命周期管理',
    '管理端 + 学生端双平台，管理员全局管控，学生自助操作',
    '实现报名→审核→考试→AI审核→归档的完整数字化闭环',
    'AI 智能审核引擎三维度自动评估，提升审核效率和标准一致性',
    'JWT 双角色认证 (admin+student)，bcrypt 密码安全存储',
    '支持 Excel 批量导入和按班级自动归档，大幅降低人工工作量',
    '公开信息门户实现证书/培训/材料的统一公示，无需登录',
    '预留教务系统对接接口，便于跨系统数据互通',
    '全栈 TypeScript + 零配置数据库 (sql.js)，开箱即用',
    '管理端 10 个功能模块 + 学生端 8 个功能模块 + 公开门户',
]

for i, ach in enumerate(achievements):
    add_text_box(slide, Inches(5.5), Inches(1.15) + Inches(0.38) * i, Inches(7.5), Inches(0.35),
                 f'✔ {ach}', font_size=10, color=DARK_TEXT)

# 统计数字
stats_data = [
    ('50+', 'API端点'),
    ('11', '数据库表'),
    ('10+8', '功能模块'),
    ('3', '证书类型'),
    ('2', '前端应用'),
    ('1', '后端服务'),
]
for i, (num, label) in enumerate(stats_data):
    left = Inches(5.5) + Inches(2.4) * i
    top = Inches(5.1)
    shape = add_rounded_rect(slide, left, top, Inches(2.0), Inches(1.1), color=PRIMARY)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.1), Inches(1.9), Inches(0.45),
                 num, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.6), Inches(1.9), Inches(0.3),
                 label, font_size=11, color=RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER)

# 快速启动
shape = add_rect(slide, Inches(5.3), Inches(6.5), Inches(7.5), Inches(0.85), color=DARK)
add_text_box(slide, Inches(5.5), Inches(6.55), Inches(7), Inches(0.3),
             '🚀 快速启动', font_size=14, color=WHITE, bold=True)
add_multi_text(slide, Inches(5.5), Inches(6.85), Inches(7), Inches(0.5), [
    'cd server && npm run dev (端口3001)   |   cd client && npm run dev (管理端:5173)   |   cd client-student && npm run dev (学生端:5174)',
], font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC), line_spacing=1.2)


# ============================================================
# Slide 13: 结束页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
             '感谢聆听', font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
             '学生校内双证管理系统', font_size=30, color=PRIMARY)

add_rect(slide, Inches(1.5), Inches(3.6), Inches(3), Inches(0.01), color=PRIMARY)

add_multi_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(3.0), [
    '📜 人社证书 · 专业证书 · 校内引进证书 — 全流程数字化管理',
    '🔧 管理端 + 🎓 学生端 — 双平台架构，各司其职',
    '🤖 AI 智能审核引擎 — 三维度自动评估，附置信度评分',
    '🔐 JWT 双角色认证 — admin + student 安全体系',
    '📦 批量导入 & 按班级归档 — 高效数据管理',
    '🌐 公开信息门户 /public — 无需登录，透明化信息公示',
    '',
    '技术栈: React 18 + TypeScript + Ant Design 5 | Node.js + Express + SQLite (sql.js)',
    '管理端: http://localhost:5173 | 学生端: http://localhost:5174 | 后端: http://localhost:3001',
    '默认账号: admin / admin123 (管理端)  |  2024001 / 024001 (学生端)',
    '',
    'GitHub: https://github.com/wangxiaolong0907/Dual-Certificate-Management-for-Students-on-Campus',
], font_size=13, color=RGBColor(0x99, 0xBB, 0xDD), line_spacing=1.6)


# ============================================================
# 保存
# ============================================================
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    '学生校内双证管理系统_路演PPT.pptx'
)
prs.save(output_path)
print('PPT saved to: ' + output_path)
print('Total slides: ' + str(len(prs.slides)))
